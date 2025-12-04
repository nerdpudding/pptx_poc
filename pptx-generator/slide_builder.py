#!/usr/bin/env python3
"""
SlideBuilder - PowerPoint presentation builder using python-pptx

Handles creation of professional presentations with consistent styling.
Supports dynamic slide counts and multiple slide types.
"""

import os
import logging
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

# Default template path inside container
DEFAULT_TEMPLATE_PATH = "/app/templates/default.pptx"

# Color scheme - Professional blue theme
COLORS = {
    "primary": RGBColor(0x1E, 0x40, 0xAF),      # Deep blue for headers
    "secondary": RGBColor(0x3B, 0x82, 0xF6),    # Lighter blue for accents
    "text_dark": RGBColor(0x1F, 0x29, 0x37),    # Dark gray for body text
    "text_light": RGBColor(0xFF, 0xFF, 0xFF),   # White for title slides
    "accent": RGBColor(0x06, 0xB6, 0xD4),       # Cyan accent
}

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


class SlideBuilder:
    """
    Builder class for creating PowerPoint presentations.

    Supports loading from template or creating blank presentations.
    Dynamically handles any number of slides.
    """

    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize SlideBuilder.

        Args:
            template_path: Path to PPTX template file. If None or not found,
                          creates a blank presentation with default styling.
        """
        template = template_path or DEFAULT_TEMPLATE_PATH

        if template and os.path.exists(template):
            logger.info(f"Loading template from: {template}")
            self.prs = Presentation(template)
            self._has_template = True
        else:
            logger.info("Creating blank presentation (no template)")
            self.prs = Presentation()
            self.prs.slide_width = SLIDE_WIDTH
            self.prs.slide_height = SLIDE_HEIGHT
            self._has_template = False

    def add_title_slide(
        self,
        heading: str,
        subheading: Optional[str] = None
    ) -> None:
        """
        Add a title slide with centered title and optional subtitle.

        Args:
            heading: Main title text
            subheading: Optional subtitle text
        """
        logger.debug(f"Adding title slide: {heading}")

        # Use blank layout and build custom
        layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(layout)

        # Add header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            SLIDE_WIDTH, Inches(2.5)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS["primary"]
        header.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.8),
            Inches(12.333), Inches(1.2)
        )
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_p = title_tf.paragraphs[0]
        title_p.text = heading
        title_p.font.size = Pt(44)
        title_p.font.bold = True
        title_p.font.color.rgb = COLORS["text_light"]
        title_p.alignment = PP_ALIGN.CENTER

        # Add subtitle if provided
        if subheading:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.0),
                Inches(12.333), Inches(0.5)
            )
            subtitle_tf = subtitle_box.text_frame
            subtitle_p = subtitle_tf.paragraphs[0]
            subtitle_p.text = subheading
            subtitle_p.font.size = Pt(24)
            subtitle_p.font.color.rgb = COLORS["text_light"]
            subtitle_p.alignment = PP_ALIGN.CENTER

    def add_content_slide(
        self,
        heading: str,
        bullets: Optional[list[str]] = None,
        subheading: Optional[str] = None
    ) -> None:
        """
        Add a content slide with title and bullet points.

        Args:
            heading: Slide title
            bullets: List of bullet point strings
            subheading: Optional subtitle (shown below title)
        """
        logger.debug(f"Adding content slide: {heading}")

        layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(layout)

        # Add header bar (thinner for content slides)
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            SLIDE_WIDTH, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS["primary"]
        header.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.7)
        )
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = heading
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = COLORS["text_light"]

        # Calculate body start position
        body_top = Inches(1.6)

        # Add subheading if provided
        if subheading:
            sub_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(1.4),
                Inches(11.833), Inches(0.4)
            )
            sub_tf = sub_box.text_frame
            sub_p = sub_tf.paragraphs[0]
            sub_p.text = subheading
            sub_p.font.size = Pt(18)
            sub_p.font.italic = True
            sub_p.font.color.rgb = COLORS["secondary"]
            body_top = Inches(1.9)

        # Add bullet points
        if bullets:
            body_box = slide.shapes.add_textbox(
                Inches(0.75), body_top,
                Inches(11.833), Inches(5.1)
            )
            body_tf = body_box.text_frame
            body_tf.word_wrap = True

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = body_tf.paragraphs[0]
                else:
                    p = body_tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(22)
                p.font.color.rgb = COLORS["text_dark"]
                p.space_before = Pt(12)
                p.space_after = Pt(6)

    def add_summary_slide(
        self,
        heading: str,
        bullets: Optional[list[str]] = None
    ) -> None:
        """
        Add a summary/conclusion slide.

        Uses slightly different styling to visually differentiate from content.

        Args:
            heading: Slide title (e.g., "Summary", "Conclusion", "Key Takeaways")
            bullets: List of summary points
        """
        logger.debug(f"Adding summary slide: {heading}")

        layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(layout)

        # Add header bar with accent color for visual distinction
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            SLIDE_WIDTH, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLORS["secondary"]  # Lighter blue
        header.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.7)
        )
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = heading
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = COLORS["text_light"]

        # Add bullet points
        if bullets:
            body_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(1.6),
                Inches(11.833), Inches(5.4)
            )
            body_tf = body_box.text_frame
            body_tf.word_wrap = True

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = body_tf.paragraphs[0]
                else:
                    p = body_tf.add_paragraph()
                p.text = f"✓ {bullet}"  # Checkmark for summary items
                p.font.size = Pt(22)
                p.font.color.rgb = COLORS["text_dark"]
                p.space_before = Pt(14)
                p.space_after = Pt(6)

    def add_slide_by_type(
        self,
        slide_type: str,
        heading: str,
        subheading: Optional[str] = None,
        bullets: Optional[list[str]] = None
    ) -> None:
        """
        Add a slide based on type string.

        Convenience method for dynamic slide creation from JSON data.

        Args:
            slide_type: One of "title", "content", "summary"
            heading: Slide heading
            subheading: Optional subtitle
            bullets: Optional list of bullet points
        """
        slide_type_lower = slide_type.lower()

        if slide_type_lower == "title":
            self.add_title_slide(heading, subheading)
        elif slide_type_lower == "content":
            self.add_content_slide(heading, bullets, subheading)
        elif slide_type_lower == "summary":
            self.add_summary_slide(heading, bullets)
        else:
            # Default to content slide for unknown types
            logger.warning(f"Unknown slide type '{slide_type}', using content layout")
            self.add_content_slide(heading, bullets, subheading)

    @property
    def slide_count(self) -> int:
        """Return the current number of slides."""
        return len(self.prs.slides)

    def save(self, file_path: str) -> str:
        """
        Save the presentation to a file.

        Args:
            file_path: Full path for the output file

        Returns:
            The file path where the presentation was saved
        """
        # Ensure output directory exists
        output_dir = os.path.dirname(file_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)

        self.prs.save(file_path)
        logger.info(f"Presentation saved: {file_path} ({self.slide_count} slides)")
        return file_path
