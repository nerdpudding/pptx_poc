#!/usr/bin/env python3
"""
Create a professional PPTX template with styled master slides.
Run once to generate default.pptx template.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme - Professional blue theme
COLORS = {
    "primary": RGBColor(0x1E, 0x40, 0xAF),      # Deep blue for headers
    "secondary": RGBColor(0x3B, 0x82, 0xF6),    # Lighter blue for accents
    "text_dark": RGBColor(0x1F, 0x29, 0x37),    # Dark gray for body text
    "text_light": RGBColor(0xFF, 0xFF, 0xFF),   # White for title slides
    "background": RGBColor(0xFF, 0xFF, 0xFF),   # White background
    "accent": RGBColor(0x06, 0xB6, 0xD4),       # Cyan accent
}


def create_template():
    """Create the default presentation template."""
    prs = Presentation()

    # Set 16:9 widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Access slide master
    slide_master = prs.slide_master

    # We can't fully modify master layouts programmatically,
    # but we'll create sample slides that demonstrate the styling
    # The SlideBuilder will apply consistent styling when creating slides

    # Create a title slide layout demo
    title_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(title_slide_layout)

    # Add a colored header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(2.5)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS["primary"]
    header.line.fill.background()

    # Add title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.8),
        Inches(12.333), Inches(1.2)
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.text = "Presentation Title"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS["text_light"]
    title_p.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.0),
        Inches(12.333), Inches(0.5)
    )
    subtitle_tf = subtitle_box.text_frame
    subtitle_p = subtitle_tf.paragraphs[0]
    subtitle_p.text = "Subtitle goes here"
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = COLORS["text_light"]
    subtitle_p.alignment = PP_ALIGN.CENTER

    # Create a content slide demo
    content_layout = prs.slide_layouts[6]  # Blank layout
    slide2 = prs.slides.add_slide(content_layout)

    # Header bar (thinner for content slides)
    content_header = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(1.2)
    )
    content_header.fill.solid()
    content_header.fill.fore_color.rgb = COLORS["primary"]
    content_header.line.fill.background()

    # Content title
    content_title = slide2.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(0.7)
    )
    ctf = content_title.text_frame
    cp = ctf.paragraphs[0]
    cp.text = "Slide Title"
    cp.font.size = Pt(32)
    cp.font.bold = True
    cp.font.color.rgb = COLORS["text_light"]

    # Content body area
    body_box = slide2.shapes.add_textbox(
        Inches(0.75), Inches(1.6),
        Inches(11.833), Inches(5.4)
    )
    body_tf = body_box.text_frame
    body_tf.word_wrap = True

    bullets = ["First bullet point", "Second bullet point", "Third bullet point"]
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = body_tf.paragraphs[0]
        else:
            p = body_tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS["text_dark"]
        p.space_before = Pt(12)

    # Save template
    output_path = os.path.join(os.path.dirname(__file__), "default.pptx")
    prs.save(output_path)
    print(f"Template created: {output_path}")
    return output_path


if __name__ == "__main__":
    create_template()
